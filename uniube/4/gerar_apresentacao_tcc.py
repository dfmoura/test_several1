#!/usr/bin/env python3
"""
Apresentacao TCC: Relato de Experiência do Projeto Integrado
Aula-04: 8 a 10 slides, 7 a 10 min, apoio visual.
Sem traco tipografico (em dash) e sem ponto mediano.
Base: Semanas 1, 3 e 5 + evidencias.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
EVID = ROOT / "evidencias"
OUT = ROOT / "apresentacao" / "Apresentacao_TCC_Relato_OSB_Uberlandia.pptx"

NAVY = RGBColor(0x1A, 0x2B, 0x4C)
BLUE = RGBColor(0x3F, 0x5D, 0x91)
ACCENT = RGBColor(0xB0, 0x7A, 0x0C)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x5A, 0x5A, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xEE, 0xF1, 0xF5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_run(run, text, size=24, bold=False, color=DARK, font="Calibri"):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def clear_tf(shape):
    shape.text_frame.clear()
    shape.text_frame.word_wrap = True


def box(slide, left, top, width, height, fill=SOFT, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is not None:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    try:
        sh.adjustments[0] = 0.06
    except Exception:
        pass
    return sh


def rect(slide, left, top, width, height, fill=NAVY):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh


def title_block(slide, title, subtitle=None):
    rect(slide, Inches(0.55), Inches(0.38), Inches(0.1), Inches(0.5), fill=ACCENT)
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11.8), Inches(0.6))
    clear_tf(tb)
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    _set_run(r, title, size=26, bold=True, color=NAVY)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.8), Inches(0.4))
        clear_tf(sub)
        p2 = sub.text_frame.paragraphs[0]
        r2 = p2.add_run()
        _set_run(r2, subtitle, size=15, bold=False, color=GRAY)


def fit_picture(slide, path, left, top, max_w, max_h):
    from PIL import Image

    with Image.open(path) as im:
        w, h = im.size
    ratio = w / h
    box_r = max_w / max_h
    if ratio > box_r:
        pw, ph = max_w, max_w / ratio
    else:
        ph, pw = max_h, max_h * ratio
    pl = left + (max_w - pw) / 2
    pt = top + (max_h - ph) / 2
    return slide.shapes.add_picture(str(path), int(pl), int(pt), width=int(pw), height=int(ph))


def caption(slide, text, left, top, width):
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.35))
    clear_tf(tb)
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    _set_run(r, text, size=12, bold=False, color=GRAY)


def fill_lines(textbox, lines, size=15, color=DARK, bold=False, align=PP_ALIGN.LEFT, space=6):
    clear_tf(textbox)
    for i, line in enumerate(lines):
        p = textbox.text_frame.paragraphs[0] if i == 0 else textbox.text_frame.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        r = p.add_run()
        _set_run(r, line, size=size, bold=bold, color=color)


def slide_capa(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
    rect(s, 0, 0, Inches(0.28), SLIDE_H, fill=NAVY)
    rect(s, 0, Inches(6.6), SLIDE_W, Inches(0.9), fill=NAVY)

    tb = s.shapes.add_textbox(Inches(0.9), Inches(1.4), Inches(11.5), Inches(0.5))
    fill_lines(tb, ["APRESENTAÇÃO"], size=34, bold=True, color=NAVY)

    tb2 = s.shapes.add_textbox(Inches(0.9), Inches(2.05), Inches(11.5), Inches(0.9))
    fill_lines(
        tb2,
        [
            "Trabalho de Conclusão de Curso",
            "Relato de Experiência do Projeto Integrado",
        ],
        size=17,
        color=GRAY,
        space=4,
    )

    tb3 = s.shapes.add_textbox(Inches(0.9), Inches(3.3), Inches(11.5), Inches(1.5))
    fill_lines(
        tb3,
        [
            "Sistema de informação de apoio ao Observatório Social",
            "do Brasil em Uberlândia: coleta, consulta e análise",
            "de licitações públicas",
        ],
        size=20,
        bold=True,
        color=DARK,
        space=2,
    )

    tb4 = s.shapes.add_textbox(Inches(0.9), Inches(5.1), Inches(10), Inches(1.1))
    fill_lines(
        tb4,
        [
            "Diogo Ferreira Moura - RA 1030125-2",
            "Sistemas de Informação - Polo Uberlândia",
            "Universidade de Uberaba - 2026/1",
        ],
        size=15,
        color=DARK,
        space=3,
    )

    foot = s.shapes.add_textbox(Inches(0.9), Inches(6.78), Inches(11), Inches(0.4))
    fill_lines(foot, ["https://licitacoes.osbrasiluberlandia.org/"], size=14, color=WHITE)


def slide_cenario(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
    title_block(s, "Organização parceira e o problema", "Onde a ação extensionista aconteceu")

    box(s, Inches(0.55), Inches(1.45), Inches(5.9), Inches(5.3), fill=WHITE, line=SOFT)
    rect(s, Inches(0.55), Inches(1.45), Inches(0.1), Inches(5.3), fill=NAVY)
    left = s.shapes.add_textbox(Inches(0.9), Inches(1.65), Inches(5.3), Inches(4.9))
    fill_lines(
        left,
        [
            "ORGANIZAÇÃO PARCEIRA",
            "",
            "Observatório Social do Brasil",
            "de Uberlândia / MG",
            "",
            "OSC apartidária, sem fins lucrativos",
            "Acompanha licitações dos órgãos municipais",
            "Unidade local desde 2015",
            "",
            "Contato operacional:",
            "Marco Aurélio Freitas Santos",
            "",
            "Ponte institucional:",
            "Lécia Queiroz (CDL Uberlândia)",
        ],
        size=15,
        space=3,
    )

    box(s, Inches(6.8), Inches(1.45), Inches(5.9), Inches(5.3), fill=WHITE, line=SOFT)
    rect(s, Inches(6.8), Inches(1.45), Inches(0.1), Inches(5.3), fill=ACCENT)
    right = s.shapes.add_textbox(Inches(7.15), Inches(1.65), Inches(5.3), Inches(4.9))
    fill_lines(
        right,
        [
            "PROBLEMA DA ENTIDADE",
            "",
            "Não faltava dado público.",
            "O que pesava era a coleta manual.",
            "",
            "Toda semana: portal -> planilha Cronograma",
            "-> só depois analisar",
            "",
            "Fontes em recortes diferentes:",
            "Power BI da PMU, Comprasnet e PNCP",
            "",
            "O pedido foi claro: uma ferramenta",
            "de apoio. Não foi curso nem oficina.",
            "Foi software para a rotina deles.",
        ],
        size=15,
        space=3,
    )


def slide_proposta(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
    title_block(s, "Proposta de intervenção", "O que foi feito e para que")

    box(s, Inches(0.55), Inches(1.4), Inches(12.2), Inches(1.7), fill=NAVY)
    tb = s.shapes.add_textbox(Inches(0.85), Inches(1.55), Inches(11.6), Inches(1.4))
    clear_tf(tb)
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    _set_run(r, "OBJETIVO GERAL", size=13, bold=True, color=ACCENT)
    p2 = tb.text_frame.add_paragraph()
    r2 = p2.add_run()
    _set_run(
        r2,
        "Desenvolver, implantar e disponibilizar um sistema de informação de apoio "
        "à coleta, consulta e análise das licitações, com custo compatível com uma OSC "
        "sem equipe de TI.",
        size=16,
        bold=False,
        color=WHITE,
    )

    objs = [
        ("1", "Identificar a rotina,\na planilha e o cenário"),
        ("2", "Construir o sistema\ncom as fontes oficiais"),
        ("3", "Colocar em uso no\ndomínio da entidade"),
        ("4", "Registrar evidências\ne retorno da equipe"),
    ]
    x = 0.55
    for num, txt in objs:
        box(s, Inches(x), Inches(3.4), Inches(2.9), Inches(2.5), fill=SOFT)
        n = s.shapes.add_textbox(Inches(x + 0.2), Inches(3.55), Inches(2.5), Inches(0.55))
        fill_lines(n, [num], size=28, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        t = s.shapes.add_textbox(Inches(x + 0.2), Inches(4.25), Inches(2.5), Inches(1.4))
        fill_lines(t, [txt], size=15, color=DARK, align=PP_ALIGN.CENTER, space=2)
        x += 3.15

    foot = s.shapes.add_textbox(Inches(0.55), Inches(6.2), Inches(12.2), Inches(0.7))
    fill_lines(
        foot,
        [
            "Sistema em produção: https://licitacoes.osbrasiluberlandia.org/",
            "Ate quatro contas (1 admin e 3 de consulta).",
        ],
        size=14,
        bold=True,
        color=NAVY,
        space=2,
    )


def slide_metodologia(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
    title_block(s, "Material e métodos", "Como a ação foi conduzida (maio a agosto de 2026)")

    steps = [
        ("Convívio", "Reuniões quase\ntoda semana\nna ACIUB"),
        ("Incremental", "Cada fatia só\navançava se a\nanterior já servia"),
        ("Implantação", "VM AWS free tier\nDocker Compose\ndomínio do OSB"),
        ("Oferta formal", "ACIUB: 30 min\nTeams sede: 1h30\nsistema no ar"),
    ]
    x = 0.55
    for i, (title, body) in enumerate(steps):
        box(s, Inches(x), Inches(1.4), Inches(2.85), Inches(2.55), fill=WHITE, line=BLUE)
        t = s.shapes.add_textbox(Inches(x + 0.15), Inches(1.55), Inches(2.55), Inches(0.45))
        fill_lines(t, [title], size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        bd = s.shapes.add_textbox(Inches(x + 0.15), Inches(2.15), Inches(2.55), Inches(1.5))
        fill_lines(bd, [body], size=14, color=DARK, align=PP_ALIGN.CENTER, space=1)
        if i < 3:
            arr = s.shapes.add_textbox(Inches(x + 2.75), Inches(2.35), Inches(0.35), Inches(0.4))
            fill_lines(arr, [">"], size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        x += 3.15

    st = s.shapes.add_textbox(Inches(0.55), Inches(4.2), Inches(12), Inches(0.35))
    fill_lines(st, ["Material utilizado"], size=15, bold=True, color=NAVY)

    items = [
        ("AWS", "1 GB RAM, 2 vCPU\nregião São Paulo"),
        ("FastAPI", "API e telas\nno mesmo serviço"),
        ("SQLite", "arquivo local\ncom backup diário"),
        ("Docker", "Compose em\nprodução"),
        ("Fontes", "CSV Power BI PMU\ne API Compras.gov"),
    ]
    x = 0.55
    for title, body in items:
        box(s, Inches(x), Inches(4.65), Inches(2.35), Inches(2.15), fill=SOFT)
        t = s.shapes.add_textbox(Inches(x + 0.1), Inches(4.8), Inches(2.15), Inches(0.4))
        fill_lines(t, [title], size=15, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        bd = s.shapes.add_textbox(Inches(x + 0.1), Inches(5.3), Inches(2.15), Inches(1.2))
        fill_lines(bd, [body], size=13, color=DARK, align=PP_ALIGN.CENTER, space=1)
        x += 2.5


def slide_arquitetura(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
    title_block(s, "Arquitetura e fluxo", "Desenho do sistema (evidência técnica)")

    left_img = EVID / "05_arquitetura.png"
    right_img = EVID / "04_fluxo.png"

    box(s, Inches(0.4), Inches(1.3), Inches(6.2), Inches(5.4), fill=SOFT)
    if left_img.exists():
        fit_picture(s, left_img, Inches(0.55), Inches(1.4), Inches(5.9), Inches(4.85))
    caption(s, "Arquitetura do sistema de apoio", Inches(0.4), Inches(6.8), Inches(6.2))

    box(s, Inches(6.8), Inches(1.3), Inches(6.1), Inches(5.4), fill=SOFT)
    if right_img.exists():
        fit_picture(s, right_img, Inches(6.95), Inches(1.4), Inches(5.8), Inches(4.85))
    caption(s, "Fluxo de coleta e consulta", Inches(6.8), Inches(6.8), Inches(6.1))


def slide_evidencias(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
    rect(s, 0, 0, SLIDE_W, Inches(0.1), fill=ACCENT)
    title_block(s, "Evidências do sistema em uso", "Prints das telas em produção")

    imgs = [
        (EVID / "01_tela_login.png", "Tela de login"),
        (EVID / "02_tela_mapa_homologados.png", "Mapa de homologados"),
        (EVID / "07_tela_cnpjs_vencedores.png", "CNPJs vencedores"),
    ]
    x = 0.4
    for path, label in imgs:
        box(s, Inches(x), Inches(1.35), Inches(4.05), Inches(5.25), fill=SOFT)
        if path.exists():
            fit_picture(s, path, Inches(x + 0.12), Inches(1.5), Inches(3.8), Inches(4.55))
        caption(s, label, Inches(x), Inches(6.7), Inches(4.05))
        x += 4.25


def slide_impacto(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
    rect(s, 0, 0, SLIDE_W, Inches(0.1), fill=ACCENT)
    title_block(s, "Resultados e retorno da entidade", "Volume carregado e o que a equipe pediu")

    metrics = [
        ("1.215", "processos municipais\nno recorte de 2025"),
        ("+36 mil", "linhas de gestores\ne fiscais"),
        ("4", "contas ativas\n(1 admin + 3)"),
        ("2", "apresentações\nformais"),
    ]
    x = 0.55
    for val, label in metrics:
        box(s, Inches(x), Inches(1.4), Inches(2.95), Inches(2.35), fill=NAVY)
        t = s.shapes.add_textbox(Inches(x + 0.15), Inches(1.55), Inches(2.65), Inches(0.8))
        fill_lines(t, [val], size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        lb = s.shapes.add_textbox(Inches(x + 0.15), Inches(2.45), Inches(2.65), Inches(1.0))
        fill_lines(lb, [label], size=13, color=WHITE, align=PP_ALIGN.CENTER, space=1)
        x += 3.15

    box(s, Inches(0.55), Inches(4.05), Inches(6.0), Inches(2.75), fill=SOFT)
    t = s.shapes.add_textbox(Inches(0.75), Inches(4.2), Inches(5.6), Inches(2.4))
    fill_lines(
        t,
        [
            '"Seria possível extrairmos estatísticas',
            'de participação... histórias que empresas',
            'participam para favorecer uma e vão',
            'fazendo rodízio nas vitórias."',
            "",
            "Marco Aurélio (OSB), WhatsApp, 11/06/2026",
        ],
        size=13,
        color=DARK,
        space=1,
    )

    box(s, Inches(6.8), Inches(4.05), Inches(6.0), Inches(2.75), fill=SOFT)
    t = s.shapes.add_textbox(Inches(7.0), Inches(4.2), Inches(5.6), Inches(2.4))
    fill_lines(
        t,
        [
            '"Que legal! O Marco tem me informado,',
            'ele está muito contente com os',
            'resultados!"',
            "",
            "Lécia Queiroz (CDL), WhatsApp, 15/07/2026",
            "",
            "Ganho indireto: a equipe passou a saber",
            "o que é uma API e que ela já existia.",
        ],
        size=13,
        color=DARK,
        space=1,
    )


def slide_conclusao(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
    title_block(s, "Conclusão", "O que ficou, o que limita e o que segue")

    blocks = [
        (
            "O que ficou",
            [
                "Sistema no ar, no domínio deles",
                "Consulta sem começar pela Cronograma",
                "Equipe entendeu o que é API",
                "Pediram ajuste (filtro por porte)",
                "Convite a voluntariado",
            ],
            NAVY,
        ),
        (
            "Limites",
            [
                "Free tier da AWS ate jan/2027",
                "Se a API Compras.gov para,",
                "só esse módulo deixa de atualizar",
                "Estudo de rodízio ainda parcial",
                "Frente dos vereadores ficou de fora",
            ],
            ACCENT,
        ),
        (
            "Próximos passos",
            [
                "Formalizar voluntariado depois",
                "da faculdade",
                "Usar menos a planilha Cronograma",
                "Medir o tempo do quadrimestre",
                "Decidir hospedagem após 2027",
            ],
            BLUE,
        ),
    ]
    x = 0.55
    for title, lines, accent in blocks:
        box(s, Inches(x), Inches(1.4), Inches(3.95), Inches(5.35), fill=WHITE, line=SOFT)
        rect(s, Inches(x), Inches(1.4), Inches(0.1), Inches(5.35), fill=accent)
        tb = s.shapes.add_textbox(Inches(x + 0.3), Inches(1.6), Inches(3.45), Inches(4.9))
        fill_lines(tb, [title, ""] + lines, size=14, color=DARK, space=4)
        # bold first line by rewriting first run
        p0 = tb.text_frame.paragraphs[0]
        if p0.runs:
            p0.runs[0].font.bold = True
            p0.runs[0].font.size = Pt(16)
            p0.runs[0].font.color.rgb = NAVY
        x += 4.2


def slide_refs_obrigado(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
    title_block(s, "Referências e encerramento", "Seleção das obras usadas no relatório")

    refs = [
        "GIL, A. C. Como elaborar projetos de pesquisa. 7. ed. Atlas, 2022.",
        "LAUDON, K.; LAUDON, J. Sistemas de informação gerenciais. 11. ed. Pearson, 2014.",
        "MOREIRA et al. Qualidade na recuperação de dados governamentais. Perspec. Ciênc. Inf., 2020.",
        "PINHO; SACRAMENTO. Accountability: já podemos traduzi-la? RAP, 2009.",
        "PRESSMAN; MAXIM. Engenharia de software. 9. ed. AMGH, 2021.",
        "SILVA. Potencial de reúso de dados abertos (DGABr). Em Questão, 2024.",
        "YIN, R. K. Estudo de caso. 5. ed. Bookman, 2015.",
        "ABNT NBR 14724, 6023 e 10520.",
    ]
    box(s, Inches(0.55), Inches(1.35), Inches(7.5), Inches(5.4), fill=SOFT)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(7.0), Inches(5.0))
    clear_tf(tb)
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    _set_run(r, "Referências (seleção)", size=14, bold=True, color=NAVY)
    for ref in refs:
        p = tb.text_frame.add_paragraph()
        p.space_before = Pt(5)
        r = p.add_run()
        _set_run(r, "- " + ref, size=11, bold=False, color=DARK)

    box(s, Inches(8.3), Inches(1.35), Inches(4.5), Inches(5.4), fill=NAVY)
    tb = s.shapes.add_textbox(Inches(8.55), Inches(2.4), Inches(4.0), Inches(3.2))
    clear_tf(tb)
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    _set_run(r, "Obrigado", size=30, bold=True, color=WHITE)
    p2 = tb.text_frame.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    _set_run(r2, "\nDiogo Ferreira Moura\nRA 1030125-2\n", size=14, bold=False, color=WHITE)
    p3 = tb.text_frame.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    _set_run(r3, "licitacoes.osbrasiluberlandia.org", size=12, bold=False, color=ACCENT)
    p4 = tb.text_frame.add_paragraph()
    p4.alignment = PP_ALIGN.CENTER
    r4 = p4.add_run()
    _set_run(r4, "\nDúvidas?", size=16, bold=True, color=WHITE)


def assert_no_forbidden(prs):
    """Garante que o PPT nao tenha em dash nem ponto mediano."""
    bad = []
    forbidden = ("\u2014", "\u2013", "\u00b7", "\u2022", "\u2192")  # em, en, middot, bullet, arrow
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            text = sh.text_frame.text
            for ch in forbidden:
                if ch in text:
                    bad.append((i, repr(ch), text[:80]))
    if bad:
        raise SystemExit("Caracteres proibidos encontrados: " + str(bad[:5]))


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_capa(prs)
    slide_cenario(prs)
    slide_proposta(prs)
    slide_metodologia(prs)
    slide_arquitetura(prs)
    slide_evidencias(prs)
    slide_impacto(prs)
    slide_conclusao(prs)
    slide_refs_obrigado(prs)

    assert_no_forbidden(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Gerado: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
